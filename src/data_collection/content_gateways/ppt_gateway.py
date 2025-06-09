import os

from pptx import Presentation

from data_collection.content_gateways import content_gateway as gateway


class PPTGateway(gateway.EducationalContentGateway):
    """
    Concrete content extractor for PowerPoint files using python-pptx.

    This reader supports extracting text from individual slides (segments).
    """

    def __init__(self):
        """
        Initializes the PPTGateway with an optional file path.
        """
        self._file_path = ""


    def get(self) -> dict[str,str]:
        """
        Extracts metadata from a PowerPoint presentation.
        """
        presentation = Presentation(self._file_path)
        props = presentation.core_properties

        file_name = os.path.basename(self._file_path)
        file_base, file_extension = os.path.splitext(file_name)
        file_extension = file_extension.lstrip(".")
        page_count = str(len(presentation.slides))

        # Extract text per slide
        pages = {}
        for i, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text)
            pages[str(i)] = "\n".join(slide_text).strip()

                # Helper to safely convert datetime to string (ISO format)
        def format_datetime(dt):
            return dt.isoformat() if dt else None
        
        metadata = {
            "file_name": file_base,
            "file_extension": file_extension,
            "page_count": page_count,
            "doc_title": props.title,
            "subject": props.subject,
            "author": props.author,
            "keywords": props.keywords,
            "comments": props.comments,
            "last_modified_by": props.last_modified_by,
            "revision": str(props.revision),
            "created": format_datetime(props.created),
            "modified": format_datetime(props.modified),
            "category": props.category,
            "content_status": props.content_status,
            "identifier": props.identifier,
            "language": props.language,
            "version": props.version,
            "pages": pages
        }
        return metadata
    
    def set_file_path(self, file_path: str) -> None:
        """
        Sets the file path for the PowerPoint file to be processed.

        Args:
            file_path (str): The path to the PowerPoint file.
        """
        self._file_path = file_path

    def get_file_path(self) -> str:
        """
        Returns the current file path set for the PowerPoint file.

        Returns:
            str: The path to the PowerPoint file.
        """
        return self._file_path


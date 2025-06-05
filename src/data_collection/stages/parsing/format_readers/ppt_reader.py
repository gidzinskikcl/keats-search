from pptx import Presentation
from data_collection.stages.parsing.format_readers import reader

class PPTReader(reader.Reader):
    """
    Concrete content extractor for PowerPoint files using python-pptx.

    This reader supports extracting text from individual slides (segments).
    """

    def get_text(self, file_path: str, segment_nr: int) -> str:
        """
        Extract text from a specific slide (segment) of a PowerPoint file.

        Args:
            file_path (str): The path to the PowerPoint (.pptx) file.
            segment_nr (int): The 1-based index of the slide to extract text from.

        Returns:
            str: The extracted text from the specified slide.

        Raises:
            IndexError: If the provided segment number is out of range.
            Exception: If an error occurs while reading the PowerPoint file.
        """
        try:
            presentation = Presentation(file_path)
            num_slides = len(presentation.slides)

            if segment_nr < 1 or segment_nr > num_slides:
                raise IndexError(f"Segment number {segment_nr} out of range. This file has {num_slides} slides.")

            slide = presentation.slides[segment_nr - 1]
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            result = "\n".join(slide_text)
        except Exception as e:
            print(f"Error reading PowerPoint file: {e}")
            raise e

        return result

    def get_images(self, file_path: str, segment_nr: int) -> list:
        """
        Extract images from a specific slide (segment) of a PowerPoint file.

        Note:
            Currently, this method is not implemented.

        Args:
            file_path (str): The path to the PowerPoint (.pptx) file.
            segment_nr (int): The 1-based index of the slide to extract images from.

        Returns:
            list: An empty list (always raises NotImplementedError).

        Raises:
            NotImplementedError: Always, since this method is not yet implemented.
        """
        raise NotImplementedError("Image extraction is not implemented yet.")

    def get_vector_graphics(self, file_path: str, segment_nr: int) -> list:
        """
        Extract vector graphics from a specific slide (segment) of a PowerPoint file.

        Note:
            Currently, this method is not implemented.

        Args:
            file_path (str): The path to the PowerPoint (.pptx) file.
            segment_nr (int): The 1-based index of the slide to extract vector graphics from.

        Returns:
            list: An empty list (always raises NotImplementedError).

        Raises:
            NotImplementedError: Always, since this method is not yet implemented.
        """
        raise NotImplementedError("Vector graphics extraction is not implemented yet.")

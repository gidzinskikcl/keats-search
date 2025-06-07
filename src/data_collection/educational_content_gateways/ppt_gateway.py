from pptx import Presentation
from data_collection.educational_content_gateways import content_gateway as gateway
from entities import segment

class PPTGateway(gateway.EducationalContentGateway):
    """
    Concrete content extractor for PowerPoint files using python-pptx.

    This reader supports extracting text from individual slides (segments).
    """

    def __init__(self):
        """
        Initializes the PPTGateway with an optional file path.
        """
        self._file_path = ""
        self._metadata = {}


    def get(self) -> list[segment.Segment]:
        """
        Extracts segments from the file by applying extraction methods.

        Args:
            file_path (str): Path to the file to be processed.
            metadata (dict[str, str]): Metadata associated with the file.
            nr_segments (int): Total number of segments to process.

        Returns:
            list[Segment]: One Segment object per segment, with aggregated text and metadata.
        """
        results = []

        nr_segments = self._metadata.get("nr_segments", 0)

        for s in range(nr_segments):
            combined_text = self.get_text(self._file_path, s + 1)
            # In the future, you might also call:
            # images = self.get_images(file_path, s + 1)
            # vector_graphics = self.get_vector_graphics(file_path, s + 1)
            # Then aggregate them as needed.

            sgmnt = segment.Segment(
                id=self._metadata.get("id", ""),
                segment_nr=s + 1,
                text=combined_text,
                file_metadata=self._metadata
            )
            results.append(sgmnt)

        return results

    @staticmethod
    def get_text(file_path: str, segment_nr: int) -> str:
        """
        Extract text from a specific slide (segment) of a PowerPoint file.

        Args:
            file_path (str): The path to the PowerPoint (.pptx) file.
            segment_nr (int): The 1-based index of the slide to extract text from.

        Returns:
            str: The extracted text from the specified slide.

        Raises:
            IndexError: If the provided segment number is out of range.
            Exception: If an error occurs while reading the PowerPoint file.
        """
        try:
            presentation = Presentation(file_path)
            num_slides = len(presentation.slides)

            if segment_nr < 1 or segment_nr > num_slides:
                raise IndexError(f"Segment number {segment_nr} out of range. This file has {num_slides} slides.")

            slide = presentation.slides[segment_nr - 1]
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            result = "\n".join(slide_text)
        except Exception as e:
            print(f"Error reading PowerPoint file: {e}")
            raise e

        return result
    
    def set_file_path(self, file_path: str) -> None:
        """
        Sets the file path for the PowerPoint file to be processed.

        Args:
            file_path (str): The path to the PowerPoint file.
        """
        self._file_path = file_path

    def get_file_path(self) -> str:
        """
        Returns the current file path set for the PowerPoint file.

        Returns:
            str: The path to the PowerPoint file.
        """
        return self._file_path
    
    def set_metadata(self, metadata: dict) -> None:
        """
        Sets the metadata for the PowerPoint file.

        Args:
            metadata (dict): Metadata associated with the PowerPoint file.
        """
        self._metadata = metadata

    def get_metadata(self) -> dict:
        """
        Returns the current metadata set for the PowerPoint file.

        Returns:
            dict: Metadata associated with the PowerPoint file.
        """
        return self._metadata
